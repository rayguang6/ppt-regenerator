from pptx import Presentation
import os
from typing import Dict, List, Any, Optional
import logging
import sys

def read_ppt(file_path: str) -> Dict[str, Any]:
    """
    Read a PowerPoint file and return detailed information about its content with tracking data.
    
    Args:
        file_path: Path to the PowerPoint file
        
    Returns:
        Dictionary containing detailed information about the presentation
    """
    # Configure logging
    logging.basicConfig(
        level=logging.INFO, 
        format='%(asctime)s - %(levelname)s: %(message)s',
        stream=sys.stdout
    )
    
    logger = logging.getLogger("PPTReader")
    logger.info(f"Opening PowerPoint file: {os.path.basename(file_path)}")
    
    prs = Presentation(file_path)
    total_slides = len(prs.slides)
    logger.info(f"Found {total_slides} slides in the presentation")
    
    # Get presentation-level information
    presentation_info = {
        "slide_count": total_slides,
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "slides": [],
        "warnings": []  # Track warnings across the presentation
    }
    
    def process_shape(shape, shape_info):
        """Helper function to process a shape and extract text"""
        # Handle grouped shapes
        if hasattr(shape, "shape_type") and shape.shape_type == 6:  # GROUP
            group_text_found = False
            if hasattr(shape, "shapes"):
                for subshape in shape.shapes:
                    subshape_info = {
                        "shape_type": str(subshape.shape_type) if hasattr(subshape, "shape_type") else "Unknown",
                        "name": subshape.name if hasattr(subshape, "name") else "Unnamed",
                        "has_text_frame": hasattr(subshape, "text_frame"),
                        "has_table": hasattr(subshape, "table"),
                        "has_chart": hasattr(subshape, "chart") and shape.shape_type == 3,
                        "text_map": []
                    }
                    process_shape(subshape, subshape_info)
                    
                    # Check if text was found in subshapes
                    if subshape_info["text_map"]:
                        group_text_found = True
                        shape_info["text_map"].extend(subshape_info["text_map"])
            
            # Add warning for grouped shape text
            if group_text_found:
                shape_info['group_text_warning'] = (
                    "Text in this grouped shape was detected but cannot be replaced. "
                    "Please ungroup shapes before processing."
                )
            
            return

        # Get text from text frames with detailed mapping
        if hasattr(shape, "text_frame"):
            try:
                text_frame = shape.text_frame
                for para_idx, paragraph in enumerate(text_frame.paragraphs):
                    if paragraph.text.strip():
                        for run_idx, run in enumerate(paragraph.runs):
                            if run.text.strip():
                                # Store detailed mapping for each run
                                run_map = {
                                    "para_idx": para_idx,
                                    "run_idx": run_idx,
                                    "text": run.text
                                }
                                shape_info["text_map"].append(run_map)
            except Exception as e:
                logger.warning(f"Could not process text frame in shape {shape.name}: {str(e)}")
    
    # Extract detailed information from each slide
    for i, slide in enumerate(prs.slides):
        # Log progress for every 5 slides or for the first/last slide
        if i == 0 or i == len(prs.slides) - 1 or (i + 1) % 5 == 0:
            logger.info(f"Processing slide {i+1} of {len(prs.slides)}")
            
        slide_info = {
            "slide_number": i + 1,
            "slide_id": slide.slide_id if hasattr(slide, "slide_id") else None,
            "slide_layout": slide.slide_layout.name if hasattr(slide.slide_layout, "name") else "Unknown",
            "shape_count": len(slide.shapes),
            "shapes": [],
            "texts": []
        }
        
        # Get information about each shape
        for j, shape in enumerate(slide.shapes):
            shape_info = {
                "shape_id": j + 1,
                "shape_type": str(shape.shape_type) if hasattr(shape, "shape_type") else "Unknown",
                "name": shape.name if hasattr(shape, "name") else "Unnamed",
                "has_text_frame": hasattr(shape, "text_frame"),
                "has_table": hasattr(shape, "table"),
                "has_chart": hasattr(shape, "chart") and shape.shape_type == 3,  # Only check chart if it's a chart shape
                "text_map": []
            }
            
            # Process the shape and its text
            process_shape(shape, shape_info)
            
            # Get table information
            if hasattr(shape, "table"):
                try:
                    table = shape.table
                    shape_info["table"] = {
                        "rows": len(table.rows),
                        "columns": len(table.columns) if len(table.rows) > 0 else 0
                    }
                except Exception as e:
                    logger.warning(f"Could not process table in shape {shape.name}: {str(e)}")
            
            slide_info["shapes"].append(shape_info)
            
            # Add texts to slide level
            for text_map in shape_info["text_map"]:
                slide_info["texts"].append(text_map["text"])
        
        presentation_info["slides"].append(slide_info)
    
    # Final analysis statistics
    text_count = sum(len(slide.get("texts", [])) for slide in presentation_info["slides"])
    shape_count = sum(len(slide.get("shapes", [])) for slide in presentation_info["slides"])
    
    logger.info(f"Analysis complete: {total_slides} slides, {shape_count} shapes, {text_count} text elements")
    
    return presentation_info

def extract_content_with_mapping(presentation_info: Dict[str, Any]) -> Dict[str, Any]:
    """
    Extract content with detailed mapping for processing.
    
    Args:
        presentation_info: Information about the presentation
        
    Returns:
        Dictionary containing extracted content with mapping
    """
    content_map = {
        "slide_count": presentation_info["slide_count"],
        "slides": [],
        "warnings": []  # Collect warnings across slides
    }
    
    for slide_idx, slide in enumerate(presentation_info["slides"]):
        slide_content = {
            "slide_number": slide["slide_number"],
            "slide_id": slide["slide_id"],
            "slide_layout": slide["slide_layout"],
            "texts": slide["texts"],
            "text_mappings": [],
            "shape_warnings": []  # Store warnings for specific shapes
        }
        
        # Create detailed mapping for each text element
        for shape_idx, shape in enumerate(slide["shapes"]):
            if shape["has_text_frame"]:
                for text_map in shape["text_map"]:
                    mapping = {
                        "shape_idx": shape_idx,
                        "para_idx": text_map["para_idx"],
                        "run_idx": text_map["run_idx"],
                        "text": text_map["text"]
                    }
                    slide_content["text_mappings"].append(mapping)
                
                # Check for group text warnings
                if "group_text_warning" in shape:
                    warning = {
                        "shape_idx": shape_idx,
                        "warning": shape["group_text_warning"]
                    }
                    slide_content["shape_warnings"].append(warning)
                    content_map["warnings"].append(warning)
        
        content_map["slides"].append(slide_content)
    
    return content_map

def modify_ppt_with_mapping(input_path: str, output_path: str, content_map: Dict[str, Any]) -> bool:
    """
    Modify PowerPoint using the content mapping to replace text while maintaining styling.
    
    Args:
        input_path: Path to the input PowerPoint file
        output_path: Path to save the modified PowerPoint file
        content_map: Mapping of content to replace
        
    Returns:
        bool: True if successful, False otherwise
    """
    try:
        prs = Presentation(input_path)
        
        # Track modifications for debugging
        modifications = []
        
        # Process each slide according to the content map
        for slide_data in content_map["slides"]:
            slide_idx = slide_data["slide_number"] - 1  # Convert to 0-based index
            
            # Skip if slide index is out of range
            if slide_idx >= len(prs.slides):
                modifications.append({
                    "error": f"Slide index {slide_idx} out of range"
                })
                continue
            
            slide = prs.slides[slide_idx]
            
            # Track text replacements for this slide
            slide_modifications = {
                "slide_number": slide_data["slide_number"],
                "replacements": []
            }
            
            # Process text mappings
            for mapping_idx, mapping in enumerate(slide_data["text_mappings"]):
                shape_idx = mapping["shape_idx"]
                para_idx = mapping["para_idx"]
                run_idx = mapping["run_idx"]
                
                # Skip if shape index is out of range
                if shape_idx >= len(slide.shapes):
                    slide_modifications["replacements"].append({
                        "error": f"Shape index {shape_idx} out of range",
                        "mapping": mapping
                    })
                    continue
                
                shape = slide.shapes[shape_idx]
                
                # Skip if shape doesn't have a text frame
                if not hasattr(shape, "text_frame"):
                    slide_modifications["replacements"].append({
                        "error": f"Shape {shape_idx} has no text frame",
                        "mapping": mapping
                    })
                    continue
                
                text_frame = shape.text_frame
                
                # Skip if paragraph index is out of range
                if para_idx >= len(text_frame.paragraphs):
                    slide_modifications["replacements"].append({
                        "error": f"Paragraph index {para_idx} out of range",
                        "mapping": mapping
                    })
                    continue
                
                paragraph = text_frame.paragraphs[para_idx]
                
                # Skip if run index is out of range
                if run_idx >= len(paragraph.runs):
                    slide_modifications["replacements"].append({
                        "error": f"Run index {run_idx} out of range",
                        "mapping": mapping
                    })
                    continue
                
                run = paragraph.runs[run_idx]
                
                # Get the new text content
                if mapping_idx < len(slide_data.get("regenerated_texts", [])):
                    new_text = slide_data["regenerated_texts"][mapping_idx]
                else:
                    # If regenerated text isn't available, use original with prefix
                    new_text = f"PLACEHOLDER: {run.text}"
                
                # Save original and new text for debugging
                replacement_info = {
                    "shape_idx": shape_idx,
                    "para_idx": para_idx,
                    "run_idx": run_idx,
                    "original_text": run.text,
                    "new_text": new_text
                }
                
                # Replace the text, preserving styling
                run.text = new_text
                
                # Record the successful replacement
                replacement_info["status"] = "success"
                slide_modifications["replacements"].append(replacement_info)
            
            modifications.append(slide_modifications)
        
        # Save the modified presentation
        prs.save(output_path)
        return True
    
    except Exception as e:
        print(f"Error modifying PowerPoint: {str(e)}")
        return False